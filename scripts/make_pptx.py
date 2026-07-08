"""Generátor prezentácie Autotest Agent (10 slajdov)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Farby (brand) ---
BG = RGBColor(0x0F, 0x17, 0x2A)        # tmavá navy
BG2 = RGBColor(0x15, 0x20, 0x38)       # panel
ACCENT = RGBColor(0x00, 0x9C, 0xC4)    # tyrkysová
ACCENT2 = RGBColor(0x3A, 0xD0, 0xB0)   # zelenkastá
WHITE = RGBColor(0xF2, 0xF6, 0xFB)
MUTED = RGBColor(0x9F, 0xB2, 0xC9)
GOLD = RGBColor(0xE6, 0xB4, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def rect(s, x, y, w, h, color, line=None):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, color, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def header(s, kicker, title):
    rect(s, Inches(0.6), Inches(0.55), Inches(0.14), Inches(0.95), ACCENT)
    txt(s, Inches(0.9), Inches(0.5), Inches(11.8), Inches(0.5),
        [[(kicker, 13, ACCENT, True)]])
    txt(s, Inches(0.9), Inches(0.85), Inches(11.8), Inches(0.9),
        [[(title, 30, WHITE, True)]])


def bullet(items, size=16):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, sub = it
            paras.append([("▸  ", size, ACCENT, True), (head, size, WHITE, True),
                          (sub, size, MUTED, False)])
        else:
            paras.append([("▸  ", size, ACCENT, True), (it, size, WHITE, False)])
    return paras


def card(s, x, y, w, h, title, lines, accent=ACCENT):
    rect(s, x, y, w, h, BG2)
    rect(s, x, y, w, Inches(0.09), accent)
    txt(s, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), Inches(0.5),
        [[(title, 16, WHITE, True)]])
    body = [[(l, 12.5, MUTED, False)] for l in lines]
    txt(s, x + Inches(0.25), y + Inches(0.72), w - Inches(0.5), h - Inches(0.9),
        body, space_after=4)


def footer(s, n):
    txt(s, Inches(0.9), Inches(7.02), Inches(9), Inches(0.35),
        [[("Autotest Agent  ·  v0.7.1", 10, MUTED, False)]])
    txt(s, Inches(11.6), Inches(7.02), Inches(1.2), Inches(0.35),
        [[(f"{n} / 10", 10, MUTED, False)]], align=PP_ALIGN.RIGHT)


# =================================================================== SLIDE 1
s = slide()
rect(s, 0, 0, SW, Inches(0.16), ACCENT)
rect(s, 0, Inches(7.34), SW, Inches(0.16), ACCENT2)
txt(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.5),
    [[("VS CODE  ·  GITHUB COPILOT  ·  MCP", 15, ACCENT, True)]])
txt(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.4),
    [[("Autotest Agent", 60, WHITE, True)]])
txt(s, Inches(1.0), Inches(3.9), Inches(11.0), Inches(1.2),
    [[("AI agent, ktorý overí opravu bugu alebo test scenár tak, že ",
       21, MUTED, False),
      ("priamo ovláda aplikáciu", 21, ACCENT2, True),
      (" — web cez Playwright MCP, desktop cez Terminator MCP.", 21, MUTED, False)]],
    line_spacing=1.15)
txt(s, Inches(1.0), Inches(5.35), Inches(11.0), Inches(0.6),
    [[("Žiadne generované test skripty. Žiadne ladenie selektorov. "
       "Scenár reálne odklikne a vizuálne overí Copilot agent mode.",
       15, WHITE, False)]], line_spacing=1.2)
footer(s, 1)

# =================================================================== SLIDE 2
s = slide()
header(s, "PROBLÉM", "Prečo klasické automatizované testy bolia")
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.2),
    bullet([
        ("Krehké skripty. ", "Selektory sa menia, testy padajú aj keď appka funguje."),
        ("Drahá údržba. ", "Viac času ide na opravu testov než na samotné testovanie."),
        ("Pomalé písanie. ", "Každý scenár treba nakódiť, odladiť a udržiavať."),
        ("Bez vizuálnej kontroly. ", "Assert na DOM nezachytí, že UI vyzerá zle."),
        ("Ručné pretestovanie bugov. ", "Tester klikne scenár znova a znova manuálne."),
    ], size=17), )
rect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.0), BG2)
txt(s, Inches(1.2), Inches(5.83), Inches(11.0), Inches(0.8),
    [[("Cieľ: ", 16, ACCENT, True),
      ("ušetriť čas testerom aj programátorom — namiesto písania a opravovania "
       "krehkých testov necháš agenta scenár reálne odklikať a vizuálne overiť.",
       16, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, 2)

# =================================================================== SLIDE 3
s = slide()
header(s, "RIEŠENIE", "Agent, ktorý appku naozaj používa")
txt(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.9),
    [[("Tester alebo dev zadá scenár (manuálne alebo z TFS bugu). Agent ho ",
       17, MUTED, False),
      ("vykoná v Copilot agent mode cez MCP nástroje", 17, ACCENT2, True),
      (" a zapíše jediný výsledok so screenshotmi krok-po-kroku.", 17, MUTED, False)]],
    line_spacing=1.2)
y = Inches(2.95)
w = Inches(3.75)
gap = Inches(0.1)
x0 = Inches(0.9)
card(s, x0, y, w, Inches(2.9), "Zadaj scenár",
     ["Manuálny popis  → test_NNN", "TFS bug  → bug_<id>",
      "LLM vygeneruje", "test_scenario.md", "(editovateľný)"], ACCENT)
card(s, x0 + w + gap, y, w, Inches(2.9), "Agent vykoná",
     ["Playwright / Terminator MCP", "klika, píše, screenshotuje",
      "auto-approve nástrojov", "pýta sa len na login /", "chýbajúci údaj"], ACCENT2)
card(s, x0 + 2 * (w + gap), y, w, Inches(2.9), "Jeden výsledok",
     ["result.md → VERDIKT", "PASSED / FAILED + zhrnutie",
      "steps/ screenshoty", "transcript.md akcií", "Dashboard + report"], GOLD)
footer(s, 3)

# =================================================================== SLIDE 4
s = slide()
header(s, "VLASTNOSTI", "Čo Autotest Agent ponúka")
feats = [
    ("🧭  Dashboard so sprievodcom",
     ["Inicializácia v krokoch: aplikácia →", "prihlásenie → TFS. Prehľad testov a reportov."]),
    ("🌐  Web  +  🖥️  Desktop",
     ["Jednotný tok cez MCP servery", "Playwright (web) a Terminator (desktop)."]),
    ("🤖  Copilot agent mode",
     ["Agent klika, píše a screenshotuje sám.", "Pýta sa len na login alebo chýbajúci údaj."]),
    ("🔗  TFS / Azure DevOps",
     ["Načíta pridelené work items a vytvorí", "test priamo z bugu."]),
    ("📊  Pekný report",
     ["Verdikt PASSED/FAILED + zhrnutie", "+ screenshot každého kroku."]),
    ("♻️  Auto-refresh",
     ["Dashboard sa po dokončení testu", "sám obnoví; indikátor „beží…\"."]),
]
w = Inches(3.75)
h = Inches(1.95)
gx = Inches(0.1)
gy = Inches(0.18)
x0 = Inches(0.9)
y0 = Inches(1.85)
for i, (t, lines) in enumerate(feats):
    cx = x0 + (i % 3) * (w + gx)
    cy = y0 + (i // 3) * (h + gy)
    card(s, cx, cy, w, h, t, lines, ACCENT if i % 2 == 0 else ACCENT2)
footer(s, 4)

# =================================================================== SLIDE 5
s = slide()
header(s, "AKO TO FUNGUJE", "Tok od scenára po report")
steps = [
    ("1", "Používateľ", "Dashboard alebo @autotest"),
    ("2", "Scenár", "test_scenario.md (LLM z popisu)"),
    ("3", "Delegácia", "Copilot agent mode + mcp.json"),
    ("4", "MCP server", "Playwright / Terminator"),
    ("5", "Vykonanie", "agent klika v aplikácii, steps/*.png"),
    ("6", "Výsledok", "result.md + transcript.md"),
    ("7", "Dashboard", "verdikt, history, report panel"),
]
y = Inches(2.1)
h = Inches(0.62)
for i, (n, t, d) in enumerate(steps):
    cy = y + i * (h + Inches(0.06))
    rect(s, Inches(0.9), cy, Inches(0.62), h, ACCENT)
    txt(s, Inches(0.9), cy, Inches(0.62), h, [[(n, 20, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.6), cy, Inches(10.8), h, BG2)
    txt(s, Inches(1.85), cy, Inches(3.0), h, [[(t, 15, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.0), cy, Inches(7.2), h, [[(d, 14, MUTED, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# =================================================================== SLIDE 6
s = slide()
header(s, "ARCHITEKTÚRA", "Tenké, jasne oddelené moduly")
mods = [
    ("extension.ts", "activate, chat participant + commands, tenký dispatcher"),
    ("setup.ts", "init wizard (QuickPick), menu nastavení"),
    ("runner.ts", "generateScenario + delegateToAgentMode"),
    ("mcp.ts", "zápis mcp.json + settings.json (auto-approve)"),
    ("dashboard.ts", "webview view + report panel + watchers"),
    ("tfs-client.ts", "import bugov z TFS / Azure DevOps"),
    ("config.ts", "konfigurácia projektu"),
    ("report.ts / util.ts", "render reportu, pomocné funkcie"),
]
w = Inches(5.75)
h = Inches(0.92)
for i, (name, desc) in enumerate(mods):
    cx = Inches(0.9) + (i % 2) * (w + Inches(0.15))
    cy = Inches(1.95) + (i // 2) * (h + Inches(0.12))
    rect(s, cx, cy, w, h, BG2)
    rect(s, cx, cy, Inches(0.09), h, ACCENT2)
    txt(s, cx + Inches(0.25), cy + Inches(0.12), w - Inches(0.4), Inches(0.4),
        [[(name, 15, ACCENT, True)]])
    txt(s, cx + Inches(0.25), cy + Inches(0.48), w - Inches(0.4), Inches(0.4),
        [[(desc, 12.5, MUTED, False)]])
footer(s, 6)

# =================================================================== SLIDE 7
s = slide()
header(s, "PLATFORMY", "Web aj desktop cez jeden tok")
card(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(3.9), "🌐  Web — Playwright MCP",
     ["", "Ovládanie prehliadača cez Playwright MCP",
      "od Microsoftu.", "",
      "• klikanie, vypĺňanie formulárov, navigácia",
      "• screenshot každého kroku",
      "• viditeľný alebo headless režim (@autotest debug)",
      "• výstupy v autotest/_mcp_output"], ACCENT)
card(s, Inches(6.7), Inches(2.0), Inches(5.6), Inches(3.9),
     "🖥️  Desktop — Terminator MCP",
     ["", "Ovládanie natívnej desktop aplikácie",
      "cez Terminator MCP.", "",
      "• UI automation nad Windows aplikáciami",
      "• rovnaký scenár, rovnaký report",
      "• žiadny pywinauto, žiadny code-gen",
      "• jednotný delegateToAgentMode tok"], ACCENT2)
txt(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.7),
    [[("MCP servery sa sťahujú automaticky cez ", 14, MUTED, False),
      ("npx", 14, ACCENT2, True),
      (" pri prvom spustení testu — netreba nič inštalovať ručne.",
       14, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# =================================================================== SLIDE 8
s = slide()
header(s, "INTEGRÁCIA", "TFS / Azure DevOps")
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.8),
    bullet([
        ("Import bugov. ", "Sekcia „TFS bugy\" načíta pridelené work items (stavy Proposed, Active)."),
        ("Test z bugu jedným klikom. ", "Scenár sa vygeneruje z popisu bugu → bug_<id>."),
        ("Zvýraznenie hotových. ", "Bug, z ktorého už test existuje, je zlatý + tlačidlo „K testu →\"."),
        ("Upozornenie na neaktuálnosť. ", "Ak sa bug zmenil, ponúkne regeneráciu scenára."),
    ], size=16))
rect(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.25), BG2)
rect(s, Inches(0.9), Inches(5.4), Inches(0.12), Inches(1.25), GOLD)
txt(s, Inches(1.2), Inches(5.55), Inches(11.0), Inches(1.0),
    [[("🔐  Bezpečné pripojenie: ", 15, GOLD, True),
      ("organization URL + projekt + PAT (scope Work Items → Read). "
       "Token sa ukladá do VS Code Secret Storage — nikdy nie do súboru "
       "ani do promptu agenta.", 15, WHITE, False)]], line_spacing=1.2)
footer(s, 8)

# =================================================================== SLIDE 9
s = slide()
header(s, "BEZPEČNOSŤ", "Dôvera a kontrola pri automatizácii")
cards = [
    ("🔑  Secret Storage",
     ["Heslá a PAT tokeny idú do VS Code", "Secret Storage — nie do autotest/",
      "ani do git, ani do promptu."]),
    ("🧩  Workspace-scope auto-approve",
     ["Auto-schvaľovanie nástrojov len pre", "tento workspace. Ostatné projekty",
      "zostanú s tvojím nastavením."]),
    ("🛡️  Bezpečnostný dialóg",
     ["Pri prvom spustení VS Code raz", "zobrazí bezpečnostný dialóg pred",
      "povolením agent mode."]),
    ("💻  Lokálne MCP servery",
     ["Bežia na tvojom stroji cez npx —", "používaj len dôveryhodné zdroje.",
      "Žiadne dáta neodchádzajú navyše."]),
]
w = Inches(5.75)
h = Inches(1.95)
for i, (t, lines) in enumerate(cards):
    cx = Inches(0.9) + (i % 2) * (w + Inches(0.15))
    cy = Inches(1.95) + (i // 2) * (h + Inches(0.2))
    card(s, cx, cy, w, h, t, lines, ACCENT if i % 2 == 0 else GOLD)
footer(s, 9)

# =================================================================== SLIDE 10
s = slide()
rect(s, 0, 0, SW, Inches(0.16), ACCENT2)
rect(s, 0, Inches(7.34), SW, Inches(0.16), ACCENT)
txt(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(0.5),
    [[("ZHRNUTIE", 15, ACCENT, True)]])
txt(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.0),
    [[("Testovanie, ktoré appku naozaj používa", 38, WHITE, True)]])
txt(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(2.4),
    bullet([
        "Bez krehkých skriptov — agent klika ako človek a vizuálne overuje.",
        "Web aj desktop cez jednotný MCP tok (Playwright / Terminator).",
        "Test z TFS bugu jedným klikom, jeden jasný verdikt so screenshotmi.",
        "Bezpečné: tokeny v Secret Storage, auto-approve len pre workspace.",
        "Beží vo VS Code v GitHub Copilot agent mode — nič navyše netreba.",
    ], size=17))
rect(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.95), BG2)
txt(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.95),
    [[("Vyskúšaj:  ", 17, ACCENT2, True),
      ("nainštaluj autotest-agent-0.7.1.vsix  →  otvor panel Autotest  →  "
       "Inicializovať projekt  →  + Test", 16, WHITE, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, 10)

out = "Autotest-Agent-prezentacia.pptx"
prs.save(out)
print("Uložené:", out)
